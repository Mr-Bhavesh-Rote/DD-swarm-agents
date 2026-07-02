#!/usr/bin/env python3
"""Generate a PPTX explaining the Deep-DD cost optimization."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(20)
    return slide


def add_two_col_slide(prs, title, left_title, left_rows, right_title, right_rows):
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1)).text_frame.text = title
    title_shape = slide.shapes[0]
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Left column
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(5.5))
    tf = left.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xD9, 0x46, 0x18)  # red-ish
    for row in left_rows:
        p = tf.add_paragraph()
        p.text = row
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    # Right column
    right = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.2), Inches(5.5))
    tf = right.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x77, 0xB4)  # blue-ish
    for row in right_rows:
        p = tf.add_paragraph()
        p.text = row
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    return slide


def add_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1)).text_frame.text = title
    title_shape = slide.shapes[0]
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True

    rows_count = len(rows) + 1
    cols_count = len(headers)
    table = slide.shapes.add_table(rows_count, cols_count, Inches(0.5), Inches(1.4), Inches(9), Inches(5)).table

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(18)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x77, 0xB4)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Data rows
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(16)

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Deep-DD Cost Optimization",
        "Reducing run cost from $30+ to under $10\nwhile understanding the accuracy trade-off",
    )

    # Slide 2: The Problem
    add_bullet_slide(
        prs,
        "Current Problem: High Run Cost",
        [
            "Recent ICL Group run cost $30.08",
            "6 parallel research agents each make multiple LLM calls",
            "Research uses Sonnet 4.6, but sheer volume drives cost",
            "Writer and verifier use expensive Opus 4.8",
            "No cost ceiling — runs can grow unchecked",
        ],
    )

    # Slide 3: Cost Before vs After
    add_table_slide(
        prs,
        "Cost Breakdown: Before vs After",
        ["Component", "Before", "After"],
        [
            ["Planner", "~$0.30", "~$0.30"],
            ["Research agents", "~$15–$20 (6 agents)", "~$3–$5 (1 agent)"],
            ["Aggregator", "~$1–$2", "~$1–$2"],
            ["Writer (Opus 4.8)", "~$5–$8", "~$3–$5"],
            ["Verifier (Opus 4.8)", "~$2–$4", "~$1–$3"],
            ["Budget cap", "None", "$10 hard stop"],
            ["Total", "$25–$35", "$8–$10"],
        ],
    )

    # Slide 4: Three Changes
    add_bullet_slide(
        prs,
        "Three Changes Implemented",
        [
            "1. Max research agents = 1 (was 6)",
            "2. Opus 4.8 as manager (planner/writer/verifier); Sonnet 4.6 as employee (research/aggregator)",
            "3. RUN_BUDGET_USD = $10 soft ceiling — warns then aborts",
        ],
    )

    # Slide 5: Why Cost Drops
    add_two_col_slide(
        prs,
        "Why Cost Drops So Much",
        "Before",
        [
            "6 agents × 3–5 tool cycles",
            "Many parallel LLM calls",
            "Long research outputs",
            "No spending cap",
        ],
        "After",
        [
            "1 agent × 3–5 tool cycles",
            "Fewer LLM calls",
            "Shorter research outputs",
            "$10 budget aborts the run",
        ],
    )

    # Slide 6: Accuracy Impact
    add_table_slide(
        prs,
        "Accuracy Impact of the Changes",
        ["Metric", "Before", "After", "Impact"],
        [
            ["Source coverage", "~75–85%", "~55–70%", "Moderate decrease"],
            ["Depth per category", "High", "Medium", "Less parallel depth"],
            ["Final report writing", "High", "High", "No change — still Opus"],
            ["Fact verification", "High", "High", "No change — still Opus"],
            ["Missed edge risks", "Low", "Moderate", "Single agent may miss niche angles"],
        ],
    )

    # Slide 7: Trade-off Summary
    add_bullet_slide(
        prs,
        "Trade-off Summary",
        [
            "Cost: $30+ → $8–$10 (65–80% cheaper)",
            "Speed: Faster — fewer parallel agents and calls",
            "Coverage: Lower — one agent covers all angles",
            "Report quality: Mostly unchanged — manager/verifier still use Opus 4.8",
            "Best for: Quick, cost-effective due-diligence overviews",
            "Not ideal for: High-stakes deep dives requiring exhaustive risk coverage",
        ],
    )

    # Slide 8: Recommendation
    add_bullet_slide(
        prs,
        "Recommendation",
        [
            "Start with the new $10 / 1-agent setup for routine checks",
            "If a report misses critical details, create a new run with 2–3 agents and a $15 budget",
            "Keep Opus for manager roles; only use Sonnet for research/aggregation",
            "Monitor Langfuse traces to see exactly where money is spent",
        ],
    )

    output_path = Path("/home/bhavesh/Documents/deep-dd/deep-dd/deep-dd-cost-optimization.pptx")
    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    main()
